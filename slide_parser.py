"""
slide_parser.py
===============
PPTXファイルを解析し、各スライドの要素（テキスト・図形・画像）を
座標情報とともに抽出するモジュール。
画像・図表はOllama上のMoondream2で内容解析を行う。

出力形式 (JSON):
{
  "file": "xxx.pptx",
  "slide_width_pt": 720.0,
  "slide_height_pt": 540.0,
  "slides": [
    {
      "slide_index": 0,
      "slide_number": 1,
      "elements": [
        {
          "id": "s0_e0",
          "type": "text",          # text / image / shape / table / smartart
          "content": "タイトル",
          "bbox_pt": [left, top, width, height],   # ポイント単位
          "bbox_ratio": [x, y, w, h],              # スライドサイズに対する比率 (0~1)
          "shape_name": "Title 1"
        },
        {
          "id": "s0_e1",
          "type": "image",
          "content": null,
          "vlm_description": "A bar chart showing...",  # Moondream2による解析結果
          "bbox_pt": [...],
          "bbox_ratio": [...],
          "shape_name": "Picture 3"
        }
      ]
    }
  ]
}

使い方:
    python slide_parser.py your_slide.pptx [--no-vlm] [--output result.json]
"""

import argparse
import base64
import io
import json
import os
import sys
from pathlib import Path
from typing import Optional

# --- 依存ライブラリ ---
try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("[ERROR] python-pptx が見つかりません。'pip install python-pptx' を実行してください。")

try:
    from PIL import Image
except ImportError:
    sys.exit("[ERROR] Pillow が見つかりません。'pip install Pillow' を実行してください。")

# Ollamaはオプション（--no-vlm 指定時は不要）
try:
    import ollama as _ollama
    OLLAMA_AVAILABLE = True
except ImportError:
    OLLAMA_AVAILABLE = False

# COMオートメーション（Windows + PowerPointがある場合のみ）
try:
    import comtypes.client
    COM_AVAILABLE = True
except ImportError:
    COM_AVAILABLE = False


# ============================================================
# 設定
# ============================================================
VLM_MODEL_CLASSIFY = "qwen3-vl:4b"    # Phase1: 図表か否かの分類
VLM_MODEL_DESCRIBE = "qwen3-vl:4b"    # Phase2: 図表の詳細説明
VLM_MODEL = VLM_MODEL_DESCRIBE         # 後方互換

VLM_PROMPT_CLASSIFY = (
    "Is this a chart, graph, or diagram? Reply only: yes or no."
)
VLM_PROMPT = (
    "Describe the visual elements in this image. "
    "If it is a bar chart, respond in this exact format:\n"
    "CHART_TYPE: vertical bar chart\n"
    "BARS:\n"
    "- label: <name>, value: <number>, x_left: <0-1>, x_right: <0-1>, y_top: <0-1>, y_bottom: <0-1>\n"
    "- label: <name>, value: <number>, x_left: <0-1>, x_right: <0-1>, y_top: <0-1>, y_bottom: <0-1>\n"
    "TITLE: <title>\n"
    "If it is a diagram or other image, describe components and relationships in plain text.\n"
    "Be concise. No thinking, no explanation."
)


# ============================================================
# ユーティリティ
# ============================================================

def emu_to_pt(emu: int) -> float:
    """EMU (English Metric Units) をポイントに変換する。1pt = 12700 EMU"""
    return emu / 12700.0


def bbox_to_ratio(left_pt, top_pt, width_pt, height_pt,
                  slide_width_pt, slide_height_pt) -> list[float]:
    """ポイント座標をスライドサイズ比率 (0~1) に変換する。"""
    return [
        round(left_pt / slide_width_pt, 4),
        round(top_pt / slide_height_pt, 4),
        round(width_pt / slide_width_pt, 4),
        round(height_pt / slide_height_pt, 4),
    ]


def pil_image_to_base64(img: "Image.Image", fmt: str = "PNG") -> str:
    """PIL Imageをbase64文字列に変換する。"""
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return base64.b64encode(buf.getvalue()).decode("utf-8")


# ============================================================
# VLM解析
# ============================================================

def analyze_image_with_vlm(img: "Image.Image", model: str = VLM_MODEL_DESCRIBE) -> str:
    """
    2段階VLM解析:
      Phase1: moondreamで「図表か否か」を判定（高速・0.2秒）
      Phase2: qwen3-vl:4bで図表の詳細説明を生成（図表のみ・約20秒）
    写真など図表でない画像はPhase1でスキップしてNoneを返す。
    最大辺を1024pxにリサイズしてGGML_ASSERTエラーを防ぐ。
    """
    if not OLLAMA_AVAILABLE:
        return "[VLM unavailable: ollama not installed]"

    # リサイズ（GGML_ASSERTエラー対策）
    MAX_SIZE = 1024
    if max(img.size) > MAX_SIZE:
        ratio = MAX_SIZE / max(img.size)
        new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
        img = img.resize(new_size, Image.LANCZOS)

    try:
        # --- Phase 1: 図表分類チェック（moondream） ---
        img_b64 = pil_image_to_base64(img)
        classify_response = _ollama.chat(
            model=VLM_MODEL_CLASSIFY,
            messages=[{
                "role": "user",
                "content": VLM_PROMPT_CLASSIFY,
                "images": [img_b64]
            }],
            options={"think": False, "num_predict": 10},
        )
        answer = classify_response["message"]["content"].strip().lower()

        # 空レスポンス・判定不能の場合は安全側（詳細解析する）に倒す
        if not answer:
            is_chart = True
            print(f"    [VLM Phase1] classify='' → 判定不能のため詳細解析へ")
        else:
            is_chart = "yes" in answer or "chart" in answer or "graph" in answer or "diagram" in answer
            print(f"    [VLM Phase1] classify={answer!r} → {'図表→詳細解析' if is_chart else 'スキップ'}")
        if not is_chart:  # ← この2行を追加
            return None
        # --- Phase 2: 詳細説明（qwen3-vl:4b） ---
        detail_response = _ollama.chat(
            model=model,
            messages=[{
                "role": "user",
                "content": VLM_PROMPT,
                "images": [img_b64]
            }],
            options={"think": False, "num_predict": 512},
        )
        raw = detail_response["message"]["content"].strip()
        if not raw:
            # thinkingから最終的な出力部分だけ抽出（thinking終了後のテキスト）
            thinking = detail_response["message"].thinking or ""
            # thinkingの最後の段落を使う
            lines = [l.strip() for l in thinking.split('\n') if l.strip()]
            raw = '\n'.join(lines[-10:]) if lines else ""
            print(f"    [VLM Phase2] thinkingから末尾抽出")
        return raw

    except Exception as e:
        return f"[VLM error: {e}]"
# ============================================================
# PowerPoint COMオートメーション：スライド画像エクスポート
# ============================================================

def export_slides_as_images(pptx_path: str,
                             output_dir: str,
                             width_px: int = 1920) -> list[str]:
    """
    PowerPoint COMオートメーションを使ってPPTXの全スライドをPNG画像として出力する。
    Windows + PowerPointがインストールされている環境でのみ動作する。

    Parameters
    ----------
    pptx_path : str
        入力PPTXファイルのパス
    output_dir : str
        画像の出力先ディレクトリ
    width_px : int
        出力画像の横幅（ピクセル）。縦はアスペクト比から自動計算。

    Returns
    -------
    list[str]
        出力された画像ファイルのパスリスト（スライド順）
    """
    if not COM_AVAILABLE:
        raise RuntimeError(
            "comtypes が見つかりません。'pip install comtypes' を実行してください。"
        )

    os.makedirs(output_dir, exist_ok=True)
    abs_pptx = os.path.abspath(pptx_path)
    abs_out  = os.path.abspath(output_dir)

    print(f"[COM] PowerPoint を起動してスライド画像を生成中...")
    ppt_app = None
    try:
        ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
        ppt_app.Visible = 1  # 1 = msoTrue

        prs = ppt_app.Presentations.Open(abs_pptx, ReadOnly=1, WithWindow=0)

        slide_w = prs.PageSetup.SlideWidth   # ポイント単位
        slide_h = prs.PageSetup.SlideHeight
        height_px = int(width_px * slide_h / slide_w)

        image_paths = []
        for i, slide in enumerate(prs.Slides):
            out_path = os.path.join(abs_out, f"slide_{i+1:03d}.png")
            slide.Export(out_path, "PNG", width_px, height_px)
            image_paths.append(out_path)
            print(f"[COM]   スライド {i+1}/{prs.Slides.Count} -> {out_path}")

        prs.Close()
        return image_paths

    finally:
        if ppt_app is not None:
            ppt_app.Quit()
            print("[COM] PowerPoint を終了しました。")


# ============================================================
# SmartArt解析
# ============================================================

# SmartArtのgraphicData URI
_SMARTART_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"

def is_smartart(shape) -> bool:
    """graphicFrame要素がSmartArtかどうかを判定する。"""
    try:
        xml_str = shape._element.xml
        return _SMARTART_URI in xml_str
    except Exception:
        return False


def extract_smartart_texts(pptx_path: str, slide_idx: int, shape_name: str) -> list[str]:
    """
    PPTXのZIPからSmartArtのdataXMLを読み、テキストノードを抽出する。
    shape_nameを使って対応するdataファイルを特定する。
    """
    import zipfile
    import re

    texts = []
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            # diagrams/data*.xml を全部読んでテキスト抽出
            data_files = sorted([f for f in z.namelist()
                                  if re.match(r'ppt/diagrams/data\d+\.xml', f)])
            # スライドのrelファイルからどのdataXMLが対応するか特定
            rel_path = f'ppt/slides/_rels/slide{slide_idx + 1}.xml.rels'
            diagram_data_file = None
            if rel_path in z.namelist():
                rels_content = z.read(rel_path).decode('utf-8')
                # diagram dataへの参照を探す
                matches = re.findall(
                    r'Target=\"\.\./diagrams/(data\d+\.xml)\"', rels_content)
                if matches:
                    diagram_data_file = f'ppt/diagrams/{matches[0]}'

            if diagram_data_file and diagram_data_file in z.namelist():
                content = z.read(diagram_data_file).decode('utf-8')
                texts = re.findall(r'<a:t[^>]*>([^<]+)</a:t>', content)
            elif data_files:
                # fallback: 全dataファイルからテキスト抽出
                for df in data_files:
                    content = z.read(df).decode('utf-8')
                    found = re.findall(r'<a:t[^>]*>([^<]+)</a:t>', content)
                    texts.extend(found)
    except Exception as e:
        texts = [f"[SmartArt parse error: {e}]"]

    return texts


def render_smartart_as_image(pptx_path: str, slide_idx: int,
                              bbox_pt: list, slide_size_pt: tuple) -> Optional["Image.Image"]:
    """
    SmartArtをスライド画像としてレンダリングし、
    bbox領域をクロップしたPIL Imageを返す。
    LibreOfficeがインストールされていれば使用、なければNoneを返す。
    """
    import subprocess
    import tempfile
    import shutil

    if shutil.which("libreoffice") is None and shutil.which("soffice") is None:
        return None

    cmd = shutil.which("libreoffice") or shutil.which("soffice")
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            # PPTXをPNG画像に変換
            subprocess.run(
                [cmd, "--headless", "--convert-to", "png",
                 "--outdir", tmpdir, pptx_path],
                capture_output=True, timeout=60
            )
            # 変換されたPNGファイルを探す
            stem = Path(pptx_path).stem
            png_files = sorted(Path(tmpdir).glob(f"{stem}*.png"))
            if slide_idx >= len(png_files):
                return None
            slide_img = Image.open(png_files[slide_idx]).convert("RGB")

            # スライド画像サイズとポイント座標のスケール計算
            img_w, img_h = slide_img.size
            sw_pt, sh_pt = slide_size_pt
            scale_x = img_w / sw_pt
            scale_y = img_h / sh_pt

            left, top, width, height = bbox_pt
            crop_box = (
                int(left * scale_x),
                int(top * scale_y),
                int((left + width) * scale_x),
                int((top + height) * scale_y),
            )
            return slide_img.crop(crop_box)
    except Exception:
        return None


# ============================================================
# 図形タイプ判定
# ============================================================

def get_shape_type_str(shape) -> str:
    """python-pptx の Shape オブジェクトから要素タイプ文字列を返す。"""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if shape.has_text_frame:
        return "text"
    if shape.has_table:
        return "table"
    # graphicFrame = SmartArt or Chart
    if is_smartart(shape):
        return "smartart"
    # グループ・コネクタ・その他図形
    return "shape"


# ============================================================
# テキスト抽出
# ============================================================

def extract_text_content(shape) -> str:
    """テキストフレームを持つ図形からテキストを抽出する。"""
    lines = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def extract_table_content(shape) -> str:
    """テーブル図形からセルテキストをCSV風に抽出する。"""
    rows = []
    for row in shape.table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(" | ".join(cells))
    return "\n".join(rows)


# ============================================================
# 画像抽出
# ============================================================

def extract_image_from_shape(shape) -> Optional["Image.Image"]:
    """PICTURE図形から PIL Image を取得する。失敗時は None を返す。"""
    try:
        image_blob = shape.image.blob
        return Image.open(io.BytesIO(image_blob)).convert("RGB")
    except Exception:
        return None


# ============================================================
# メイン解析関数
# ============================================================

def parse_pptx(pptx_path: str, use_vlm: bool = True) -> dict:
    """
    PPTXファイルを解析し、スライド要素メタデータの辞書を返す。

    Parameters
    ----------
    pptx_path : str
        入力PPTXファイルのパス
    use_vlm : bool
        True の場合、画像・SmartArt要素をMoondream2で解析する

    Returns
    -------
    dict
        スライド要素メタデータ
    """
    prs = Presentation(pptx_path)

    slide_width_pt = emu_to_pt(prs.slide_width)
    slide_height_pt = emu_to_pt(prs.slide_height)
    slide_size_pt = (slide_width_pt, slide_height_pt)

    result = {
        "file": os.path.basename(pptx_path),
        "slide_width_pt": round(slide_width_pt, 2),
        "slide_height_pt": round(slide_height_pt, 2),
        "slides": []
    }

    total_slides = len(prs.slides)
    for slide_idx, slide in enumerate(prs.slides):
        print(f"  Processing slide {slide_idx + 1}/{total_slides}...")

        slide_data = {
            "slide_index": slide_idx,
            "slide_number": slide_idx + 1,
            "elements": []
        }

        for elem_idx, shape in enumerate(slide.shapes):
            elem_id = f"s{slide_idx}_e{elem_idx}"
            shape_type = get_shape_type_str(shape)

            # --- 座標情報 ---
            left_pt = emu_to_pt(shape.left) if shape.left is not None else 0.0
            top_pt = emu_to_pt(shape.top) if shape.top is not None else 0.0
            width_pt = emu_to_pt(shape.width) if shape.width is not None else 0.0
            height_pt = emu_to_pt(shape.height) if shape.height is not None else 0.0

            bbox_pt = [
                round(left_pt, 2),
                round(top_pt, 2),
                round(width_pt, 2),
                round(height_pt, 2)
            ]
            bbox_ratio = bbox_to_ratio(
                left_pt, top_pt, width_pt, height_pt,
                slide_width_pt, slide_height_pt
            )

            elem = {
                "id": elem_id,
                "type": shape_type,
                "content": None,
                "bbox_pt": bbox_pt,
                "bbox_ratio": bbox_ratio,
                "shape_name": shape.name,
            }

            # --- コンテンツ抽出 ---
            if shape_type == "text":
                elem["content"] = extract_text_content(shape)

            elif shape_type == "table":
                elem["content"] = extract_table_content(shape)

            elif shape_type == "image":
                pil_img = extract_image_from_shape(shape)
                if pil_img and use_vlm:
                    print(f"    [VLM] Analyzing image in element {elem_id}...")
                    elem["vlm_description"] = analyze_image_with_vlm(pil_img)
                else:
                    elem["vlm_description"] = None

            elif shape_type == "smartart":
                # ① ZIPからテキストノードを抽出（必ず実行）
                texts = extract_smartart_texts(pptx_path, slide_idx, shape.name)
                elem["content"] = texts  # リスト形式で保持
                elem["smartart_node_count"] = len(texts)
                print(f"    [SmartArt] slide{slide_idx+1} '{shape.name}': "
                      f"{len(texts)}ノード → {texts}")

                # ② VLM有効時: LibreOfficeでレンダリングしてVLM解析（オプション）
                elem["vlm_description"] = None
                if use_vlm:
                    print(f"    [VLM] Attempting SmartArt rendering for {elem_id}...")
                    cropped = render_smartart_as_image(
                        pptx_path, slide_idx, bbox_pt, slide_size_pt)
                    if cropped:
                        elem["vlm_description"] = analyze_image_with_vlm(cropped)
                        print(f"    [VLM] SmartArt VLM analysis complete.")
                    else:
                        print(f"    [VLM] LibreOffice not found, skipping render.")

            elif shape_type == "shape":
                # テキストなし図形でもテキストフレームを持つ場合がある
                if shape.has_text_frame:
                    text = extract_text_content(shape)
                    if text:
                        elem["content"] = text

            slide_data["elements"].append(elem)

        result["slides"].append(slide_data)

    return result


# ============================================================
# ASR initial_prompt 生成（LLMによる要約）
# ============================================================

_ASR_PROMPT_SYSTEM = """You are a helpful assistant for a speech recognition system.
Your task is to generate a short Japanese passage (2-3 sentences, max 200 characters)
that contains the key technical terms, proper nouns, and keywords from a presentation slide.
This passage will be used as an "initial_prompt" for the Whisper speech recognition model
to improve its accuracy for domain-specific vocabulary.

Rules:
- Write in natural Japanese (not a list)
- Include ALL technical terms, proper nouns, acronyms, and numbers from the slide
- Keep it concise: 2-3 sentences, under 200 characters
- Do NOT explain the slide content, just naturally weave the terms into sentences
- Respond ONLY with the Japanese passage, no preamble or explanation
"""

def generate_asr_prompts(
    slide_data: dict,
    model: str = "llama3",
    force: bool = False,
) -> dict:
    """
    slide_data の各スライドに asr_prompt フィールドを生成・追加する。

    Parameters
    ----------
    slide_data : dict
        parse_pptx() の返り値
    model : str
        使用するOllamaモデル名
    force : bool
        True の場合、既存の asr_prompt を上書きする

    Returns
    -------
    dict
        asr_prompt フィールドが追加された slide_data（インプレース変更）
    """
    if not OLLAMA_AVAILABLE:
        print("[ASR_PROMPT] ollama が利用できないためスキップします。")
        return slide_data

    total = len(slide_data.get("slides", []))
    print(f"[ASR_PROMPT] {total}スライド分の initial_prompt を生成中...")

    for slide in slide_data["slides"]:
        slide_num = slide["slide_number"]

        # 既に生成済みでforceでなければスキップ
        if "asr_prompt" in slide and not force:
            print(f"  スライド {slide_num:2d}: スキップ（既存）")
            continue

        # スライドの全テキストを収集
        texts = []
        for e in slide.get("elements", []):
            content = e.get("content")
            vlm_desc = e.get("vlm_description")
            # タイトル要素を優先して先頭に
            is_title = ("タイトル" in e.get("shape_name", "") or
                        "Title" in e.get("shape_name", ""))
            if isinstance(content, list):
                entry = " ".join(content)
            elif isinstance(content, str) and content.strip():
                entry = content.strip()
            elif vlm_desc:
                entry = f"[図表] {vlm_desc[:80]}"
            else:
                continue

            if is_title:
                texts.insert(0, entry)
            else:
                texts.append(entry)

        if not texts:
            slide["asr_prompt"] = "プレゼンテーション スライド 発表"
            print(f"  スライド {slide_num:2d}: テキストなし → デフォルト")
            continue

        # LLMへの入力（テキストが長すぎる場合は切り詰め）
        slide_text = "\n".join(texts)
        if len(slide_text) > 800:
            slide_text = slide_text[:800] + "..."

        user_prompt = f"""以下はプレゼンテーションのスライド{slide_num}の内容です。
このスライドに登場する専門用語・固有名詞・キーワードをすべて含む
自然な日本語の文章（2〜3文、200文字以内）を生成してください。

--- スライド内容 ---
{slide_text}
---
"""
        try:
            response = _ollama.chat(
                model=model,
                messages=[
                    {"role": "system", "content": _ASR_PROMPT_SYSTEM},
                    {"role": "user",   "content": user_prompt},
                ],
                options={"num_predict": 150, "temperature": 0.3},
            )
            prompt = response["message"]["content"].strip()
            # 200文字（≒100トークン程度）に収める
            if len(prompt) > 200:
                prompt = prompt[:200]
            slide["asr_prompt"] = prompt
            print(f"  スライド {slide_num:2d}: {prompt[:60]}...")
        except Exception as e:
            # エラー時はフォールバック
            fallback = "プレゼンテーション スライド 発表 " + " ".join(texts[:3])
            slide["asr_prompt"] = fallback[:200]
            print(f"  スライド {slide_num:2d}: LLMエラー({e}) → フォールバック")

    print("[ASR_PROMPT] 生成完了")
    return slide_data


# ============================================================
# CLIエントリポイント
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="PPTXスライド解析モジュール - 要素の座標とコンテンツをJSONで出力"
    )
    parser.add_argument("pptx_file", help="入力PPTXファイルのパス")
    parser.add_argument(
        "--no-vlm",
        action="store_true",
        help="VLM（Moondream2）による画像解析をスキップする"
    )
    parser.add_argument(
        "--output", "-o",
        default=None,
        help="出力JSONファイルのパス（省略時は入力ファイル名.jsonに保存）"
    )
    parser.add_argument(
        "--pretty",
        action="store_true",
        default=True,
        help="JSONを整形して出力する（デフォルト: True）"
    )
    parser.add_argument(
        "--export-images",
        action="store_true",
        help="PowerPoint COMでスライド画像を自動エクスポートする（Windows限定）"
    )
    parser.add_argument(
        "--images-dir",
        default=None,
        help="画像出力先ディレクトリ（省略時はPPTXファイル名_imagesフォルダ）"
    )
    parser.add_argument(
        "--image-width",
        type=int,
        default=1920,
        help="出力画像の横幅px（デフォルト: 1920）"
    )
    parser.add_argument(
        "--gen-asr-prompts",
        action="store_true",
        help="LLMでスライドごとのASR initial_prompt を生成してJSONに保存する"
    )
    parser.add_argument(
        "--asr-prompt-model",
        default="llama3",
        help="ASR prompt生成に使うOllamaモデル名（デフォルト: llama3）"
    )
    parser.add_argument(
        "--force-asr-prompts",
        action="store_true",
        help="既存のasr_promptを上書きして再生成する"
    )
    args = parser.parse_args()

    pptx_path = args.pptx_file
    if not os.path.exists(pptx_path):
        sys.exit(f"[ERROR] ファイルが見つかりません: {pptx_path}")

    use_vlm = not args.no_vlm
    if use_vlm and not OLLAMA_AVAILABLE:
        print("[WARNING] ollama ライブラリが見つかりません。VLM解析をスキップします。")
        use_vlm = False

    print(f"[INFO] Parsing: {pptx_path}")
    print(f"[INFO] VLM解析: {'有効 (qwen3-vl:4b)' if use_vlm else '無効'}")

    data = parse_pptx(pptx_path, use_vlm=use_vlm)

    # ASR initial_prompt の生成（オプション）
    if args.gen_asr_prompts:
        if not OLLAMA_AVAILABLE:
            print("[WARNING] ollama が利用できないため --gen-asr-prompts をスキップします。")
        else:
            print(f"[INFO] ASR prompt生成モデル: {args.asr_prompt_model}")
            generate_asr_prompts(
                data,
                model=args.asr_prompt_model,
                force=args.force_asr_prompts,
            )

    # 出力パス決定
    if args.output:
        out_path = args.output
    else:
        out_path = str(Path(pptx_path).stem) + "_slide_data.json"

    indent = 2 if args.pretty else None
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=indent)

    print(f"[INFO] 保存完了: {out_path}")

    # 画像エクスポート（オプション）
    if args.export_images:
        images_dir = args.images_dir or (str(Path(pptx_path).stem) + "_images")
        try:
            image_paths = export_slides_as_images(
                pptx_path, images_dir, width_px=args.image_width)
            print(f"[INFO] 画像エクスポート完了: {len(image_paths)}枚 -> {images_dir}/")
            # JSONにも画像パスを埋め込む
            for slide in data["slides"]:
                idx = slide["slide_index"]
                if idx < len(image_paths):
                    slide["image_path"] = image_paths[idx]
            # 画像パス付きで再保存
            with open(out_path, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=indent)
            print(f"[INFO] JSONに画像パスを追記して再保存: {out_path}")
        except Exception as e:
            print(f"[WARNING] 画像エクスポート失敗: {e}")

    # サマリー表示
    total_elements = sum(len(s["elements"]) for s in data["slides"])
    print(f"\n=== 解析サマリー ===")
    print(f"  スライド数     : {len(data['slides'])}")
    print(f"  総要素数       : {total_elements}")
    for slide in data["slides"]:
        type_counts = {}
        for e in slide["elements"]:
            type_counts[e["type"]] = type_counts.get(e["type"], 0) + 1
        print(f"  スライド {slide['slide_number']:2d}    : {type_counts}")


if __name__ == "__main__":
    main()
