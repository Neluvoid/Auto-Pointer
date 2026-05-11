"""
vlm_comparison.py
=================
moondream vs Qwen2.5-VL:7b の VLM性能比較スクリプト（4.1 意味理解精度向上）

PPTXファイルから画像要素を直接抽出し、各VLMに同じ画像・同じプロンプトを
渡して出力を比較する。

評価軸:
  1. 説明の質（グラフの軸・凡例・傾向・数値が読み取れているか）
  2. 日本語テキストの認識（スライド内の日本語ラベルを読めるか）
  3. 速度（vlm_description生成にかかる時間）

使い方:
    python vlm_comparison.py --pptx your_slide.pptx
    python vlm_comparison.py --pptx your_slide.pptx --models moondream qwen2.5-vl:7b
    python vlm_comparison.py --pptx your_slide.pptx --slide-nums 2 3 4
"""

import argparse
import base64
import io
import json
import os
import sys
import time
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt
except ImportError:
    sys.exit("[ERROR] python-pptx が必要です: pip install python-pptx")

try:
    from PIL import Image
except ImportError:
    sys.exit("[ERROR] Pillow が必要です: pip install Pillow")

try:
    import ollama as _ollama
except ImportError:
    sys.exit("[ERROR] ollama が必要です: pip install ollama")


# ============================================================
# プロンプト定義
# ============================================================

# slide_parser.py と同じプロンプト（公平な比較のため）
VLM_PROMPT_EN = (
    "Describe all visual elements in this image in detail. "
    "If it is a chart or graph, identify the chart type, axes labels, "
    "legend items, and notable data points or trends. "
    "If it is a diagram, describe the components and their relationships. "
    "Respond in English."
)

# 日本語対応テスト用プロンプト（Qwen2.5-VLの強みを引き出す）
VLM_PROMPT_JA = (
    "この画像に含まれるすべての視覚要素を詳細に説明してください。"
    "グラフや図表の場合は、グラフの種類、軸のラベル、凡例、"
    "注目すべきデータポイントや傾向を説明してください。"
    "日本語のテキストがあればそのまま読み取ってください。"
)


# ============================================================
# PPTXから画像要素を抽出
# ============================================================

def extract_images_from_pptx(pptx_path: str,
                              slide_nums: list[int] | None = None
                              ) -> list[dict]:
    """
    PPTXから画像要素（PICTURE型）を抽出してPIL Imageとして返す。

    Returns
    -------
    list[dict]
        [{slide_num, element_id, shape_name, image: PIL.Image, bbox_ratio}]
    """
    prs = Presentation(pptx_path)
    sw_emu = prs.slide_width
    sh_emu = prs.slide_height

    results = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        if slide_nums and slide_num not in slide_nums:
            continue

        for elem_idx, shape in enumerate(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            try:
                blob = shape.image.blob
                img  = Image.open(io.BytesIO(blob)).convert("RGB")
            except Exception as e:
                print(f"  [WARN] 画像抽出失敗 slide{slide_num} shape={shape.name}: {e}")
                continue

            # bbox_ratio 計算
            left  = shape.left  or 0
            top   = shape.top   or 0
            width = shape.width or 0
            height= shape.height or 0
            bbox_ratio = [
                round(left   / sw_emu, 4),
                round(top    / sh_emu, 4),
                round(width  / sw_emu, 4),
                round(height / sh_emu, 4),
            ]

            results.append({
                "slide_num":   slide_num,
                "element_id":  f"s{slide_idx}_e{elem_idx}",
                "shape_name":  shape.name,
                "image":       img,
                "bbox_ratio":  bbox_ratio,
                "img_size":    img.size,
            })

    return results


# ============================================================
# VLM推論
# ============================================================

def pil_to_b64(img: Image.Image, fmt: str = "PNG") -> str:
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return base64.b64encode(buf.getvalue()).decode("utf-8")


def run_vlm(model: str, img: Image.Image, prompt: str) -> tuple[str, float]:
    """VLMに画像とプロンプトを渡して説明文と所要時間を返す。"""
    img_b64 = pil_to_b64(img)
    t0 = time.time()
    response = _ollama.chat(
        model=model,
        messages=[{
            "role":    "user",
            "content": prompt,
            "images":  [img_b64],
        }]
    )
    elapsed = time.time() - t0
    return response["message"]["content"].strip(), elapsed


# ============================================================
# 比較実行
# ============================================================

def run_comparison(
    pptx_path: str,
    models: list[str],
    slide_nums: list[int] | None,
    prompts: dict[str, str],
    output_dir: str,
):
    print(f"\n[INFO] PPTX: {pptx_path}")
    print(f"[INFO] モデル: {models}")

    # 画像抽出
    images = extract_images_from_pptx(pptx_path, slide_nums)
    if not images:
        print("[ERROR] 画像要素が見つかりませんでした")
        return

    print(f"[INFO] 抽出画像数: {len(images)}枚\n")

    # モデルの事前確認
    available_models = []
    for model in models:
        try:
            _ollama.chat(
                model=model,
                messages=[{"role": "user", "content": "ping"}],
                options={"num_predict": 3},
            )
            available_models.append(model)
            print(f"  ✅ {model} : 利用可能")
        except Exception as e:
            print(f"  ❌ {model} : 利用不可 ({e})")

    if not available_models:
        sys.exit("[ERROR] 利用可能なモデルがありません")

    # 結果格納
    all_results = []

    for item in images:
        slide_num  = item["slide_num"]
        elem_id    = item["element_id"]
        shape_name = item["shape_name"]
        img        = item["image"]
        w, h       = item["img_size"]

        print(f"\n{'='*65}")
        print(f"  スライド{slide_num} / {elem_id} / {shape_name}  ({w}x{h}px)")
        print(f"{'='*65}")

        # 画像を保存（目視確認用）
        os.makedirs(output_dir, exist_ok=True)
        img_save_path = os.path.join(
            output_dir, f"slide{slide_num}_{elem_id}.png")
        img.save(img_save_path)
        print(f"  画像保存: {img_save_path}")

        for prompt_label, prompt_text in prompts.items():
            print(f"\n  --- プロンプト: {prompt_label} ---")
            for model in available_models:
                try:
                    desc, elapsed = run_vlm(model, img, prompt_text)
                    print(f"\n  [{model}]  ({elapsed:.1f}秒)")
                    # 出力を整形して表示（長い場合は折り返す）
                    for line in desc.split("\n"):
                        print(f"    {line}")

                    all_results.append({
                        "slide_num":    slide_num,
                        "element_id":   elem_id,
                        "shape_name":   shape_name,
                        "model":        model,
                        "prompt_label": prompt_label,
                        "elapsed_sec":  round(elapsed, 2),
                        "description":  desc,
                    })

                except Exception as e:
                    print(f"\n  [{model}] エラー: {e}")
                    all_results.append({
                        "slide_num":    slide_num,
                        "element_id":   elem_id,
                        "shape_name":   shape_name,
                        "model":        model,
                        "prompt_label": prompt_label,
                        "elapsed_sec":  None,
                        "description":  f"[ERROR] {e}",
                    })

    # ============================================================
    # サマリー（速度比較）
    # ============================================================
    print(f"\n\n{'='*65}")
    print("  速度サマリー（モデル別平均）")
    print(f"{'='*65}")

    for model in available_models:
        times = [r["elapsed_sec"] for r in all_results
                 if r["model"] == model and r["elapsed_sec"] is not None]
        if times:
            print(f"  {model:<25} 平均: {sum(times)/len(times):.1f}秒  "
                  f"(min: {min(times):.1f}s / max: {max(times):.1f}s)")

    # ============================================================
    # JSON保存（全出力を保存して後から評価できるように）
    # ============================================================
    from datetime import datetime
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    json_path = os.path.join(output_dir, f"vlm_comparison_{ts}.json")
    with open(json_path, "w", encoding="utf-8") as f:
        # imageキーは除外して保存
        save_data = [
            {k: v for k, v in item.items() if k != "image"}
            for item in all_results
        ]
        json.dump(save_data, f, ensure_ascii=False, indent=2)
    print(f"\n  結果JSON保存: {json_path}")

    return all_results


# ============================================================
# CLIエントリポイント
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="VLMモデル比較（moondream vs Qwen2.5-VL）"
    )
    parser.add_argument("--pptx", required=True,
                        help="入力PPTXファイルのパス")
    parser.add_argument(
        "--models", nargs="+",
        default=["moondream", "qwen3-vl:4b"],
        help="比較するVLMモデル名（デフォルト: moondream qwen2.5-vl:7b）"
    )
    parser.add_argument(
        "--slide-nums", nargs="+", type=int, default=None,
        help="比較対象のスライド番号（省略時は全スライド）"
    )
    parser.add_argument(
        "--prompt", default="both",
        choices=["en", "ja", "both"],
        help="使用するプロンプト言語（デフォルト: both）"
    )
    parser.add_argument(
        "--output-dir", default="vlm_comparison_output",
        help="出力ディレクトリ（デフォルト: vlm_comparison_output）"
    )
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        sys.exit(f"[ERROR] PPTXが見つかりません: {args.pptx}")

    # プロンプト選択
    prompts = {}
    if args.prompt in ("en", "both"):
        prompts["EN"] = VLM_PROMPT_EN
    if args.prompt in ("ja", "both"):
        prompts["JA"] = VLM_PROMPT_JA

    run_comparison(
        pptx_path=args.pptx,
        models=args.models,
        slide_nums=args.slide_nums,
        prompts=prompts,
        output_dir=args.output_dir,
    )


if __name__ == "__main__":
    main()
